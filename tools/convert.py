"""
文档转换工具：将 PDF/Word/PPT 转为 TXT
扫描 subjects/ 下所有科目的 materials/ 目录，转换后输出到 docs/
"""
import os
import sys

# --- 各格式解析器 ---

def parse_pdf(filepath):
    """从 PDF 提取纯文本"""
    from PyPDF2 import PdfReader
    reader = PdfReader(filepath)
    lines = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            lines.append(f"--- 第{i+1}页 ---")
            lines.append(text.strip())
    return "\n\n".join(lines)


def parse_docx(filepath):
    """从 Word 文档提取纯文本"""
    from docx import Document
    doc = Document(filepath)
    lines = []
    for para in doc.paragraphs:
        if para.text.strip():
            lines.append(para.text.strip())
    # 也提取表格中的文本
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)


def parse_pptx(filepath):
    """从 PPT 提取纯文本"""
    from pptx import Presentation
    prs = Presentation(filepath)
    lines = []
    for i, slide in enumerate(prs.slides):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)
        if slide_lines:
            lines.append(f"--- 第{i+1}页 ---")
            lines.extend(slide_lines)
    return "\n".join(lines)


# --- 格式映射 ---

PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
}

# --- 主逻辑 ---

def convert_subject(subject_path):
    """处理单个科目目录"""
    subject_name = os.path.basename(subject_path)
    materials_dir = os.path.join(subject_path, "materials")
    docs_dir = os.path.join(subject_path, "docs")

    if not os.path.isdir(materials_dir):
        return []

    os.makedirs(docs_dir, exist_ok=True)
    results = []

    for filename in os.listdir(materials_dir):
        filepath = os.path.join(materials_dir, filename)
        if not os.path.isfile(filepath):
            continue

        ext = os.path.splitext(filename)[1].lower()
        if ext not in PARSERS:
            continue

        base = os.path.splitext(filename)[0]
        out_path = os.path.join(docs_dir, base + ".txt")

        # 不覆盖已存在的 txt（保护手动修改）
        if os.path.exists(out_path):
            results.append((filename, "跳过（已存在）"))
            continue

        try:
            print(f"  转换中: {filename} ...", end=" ", flush=True)
            content = PARSERS[ext](filepath)
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(content)
            print("OK")
            results.append((filename, "成功"))
        except Exception as e:
            print(f"失败: {e}")
            results.append((filename, f"失败: {e}"))

    return results


def main():
    root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    subjects_root = os.path.join(root, "subjects")

    if not os.path.isdir(subjects_root):
        print("错误: 未找到 subjects/ 目录")
        return

    total_results = []

    for name in sorted(os.listdir(subjects_root)):
        subject_path = os.path.join(subjects_root, name)
        if not os.path.isdir(subject_path):
            continue
        print(f"\n[{name}]")
        results = convert_subject(subject_path)
        total_results.extend([(name, r[0], r[1]) for r in results])
        if not results:
            print("  （无待转换文件）")

    # 摘要
    print("\n" + "=" * 50)
    if not total_results:
        print("没有找到需要转换的文件。")
        print("请将 PDF/Word/PPT 文件放入 subjects/<科目名>/materials/ 目录。")
    else:
        for subject, fname, status in total_results:
            print(f"  [{subject}] {fname} -> {status}")

    try:
        print("\n按任意键退出...")
        input()
    except (EOFError, OSError):
        pass


if __name__ == "__main__":
    main()
